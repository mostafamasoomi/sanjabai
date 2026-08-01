"""Document Generator endpoints.

POST /v1/documents/generate  — Generate PPTX/DOCX/MDX from AI prompt
GET  /v1/documents           — List user's generated documents
GET  /v1/documents/{id}/download — Download generated document
DELETE /v1/documents/{id}    — Delete generated document
"""
from __future__ import annotations

import os
import re
import json
import time
import uuid
import logging
from datetime import datetime, timezone
from pathlib import Path
from typing import Optional

from fastapi import APIRouter, Request
from fastapi.responses import JSONResponse, FileResponse, Response

from database import async_session, _http, LITELLM_HOST
from dependencies import _get_user_id

logger = logging.getLogger(__name__)
router = APIRouter()

# ── Storage ──────────────────────────────────────────────────────
DOC_STORAGE = Path('/tmp/sanjhubai_docs')
DOC_STORAGE.mkdir(parents=True, exist_ok=True)

# In-memory doc registry (persists until container restart)
_doc_registry: dict[str, dict] = {}

# ── Safety constants ─────────────────────────────────────────────
MAX_PROMPT_LENGTH = 4000
MAX_SLIDES = 30
MAX_SECTIONS = 30
ALLOWED_MODELS = {
    'tencent-hy3', 'mimo-v2.5-pro', 'mimo-v2.5-pro-ultraspeed',
    'deepseek-v4-pro', 'deepseek-v4-flash', 'deepseek-v4-pro-bynara',
    'mistral-large', 'mistral-medium-3.5',
}
DOC_ID_RE = re.compile(r'^[0-9a-f]{12}$')

# ── AI content generation ────────────────────────────────────────

async def _generate_content(prompt: str, doc_type: str, model: str = 'mimo-v2.5-pro') -> dict:
    """Use AI model to generate structured content for the document."""
    import httpx

    type_instructions = {
        'pptx': '''Return a JSON object with this exact structure:
{
  "title": "Presentation Title",
  "subtitle": "Optional subtitle",
  "slides": [
    {
      "title": "Slide Title",
      "content": ["Bullet point 1", "Bullet point 2", "Bullet point 3"],
      "notes": "Optional speaker notes"
    }
  ]
}
Create 6-10 slides. Make content professional and detailed.''',

        'docx': '''Return a JSON object with this exact structure:
{
  "title": "Document Title",
  "subtitle": "Optional subtitle",
  "sections": [
    {
      "heading": "Section Heading",
      "content": "Full paragraph content for this section.",
      "subsections": [
        {
          "heading": "Subsection",
          "content": "Subsection content."
        }
      ]
    }
  ]
}
Create 4-8 sections. Make content professional and detailed.''',
    }

    system_msg = f'''You are a professional document/content creator. 
Generate structured content for a {doc_type.upper()} document based on the user's request.
{type_instructions.get(doc_type, type_instructions['pptx'])}

IMPORTANT: Return ONLY valid JSON. No markdown code blocks, no explanations. Just the raw JSON.'''

    # DOCX/PPTX content is long — use a dedicated longer timeout + one retry
    max_tokens = 4096 if doc_type != 'docx' else 3072
    timeout = httpx.Timeout(120.0, connect=15.0)
    payload = {
        'model': model,
        'messages': [
            {'role': 'system', 'content': system_msg},
            {'role': 'user', 'content': prompt},
        ],
        'stream': False,
        'max_tokens': max_tokens,
    }
    last_err: Exception | None = None
    r = None
    for attempt in range(2):
        try:
            r = await _http.post(
                f'{LITELLM_HOST}/v1/chat/completions',
                json=payload,
                headers={'Accept': 'application/json'},
                timeout=timeout,
            )
            r.raise_for_status()
            break
        except (httpx.TimeoutException, httpx.HTTPStatusError, httpx.TransportError) as e:
            last_err = e
            logger.warning('docgen AI attempt %s failed: %s', attempt + 1, e)
            if attempt == 0:
                # fallback model for second try if primary is slow/unavailable
                if payload['model'] != 'mimo-v2.5-pro':
                    payload['model'] = 'mimo-v2.5-pro'
                continue
            raise
    if r is None:
        raise last_err or RuntimeError('document content generation failed')
    content = r.json()['choices'][0]['message']['content']

    # Try to parse JSON from response — handle markdown code blocks
    content = content.strip()
    if content.startswith('```'):
        # Remove ```json or ``` wrapper
        lines = content.split('\n')
        if lines[0].startswith('```'):
            lines = lines[1:]
        if lines and lines[-1].strip() == '```':
            lines = lines[:-1]
        content = '\n'.join(lines).strip()
    
    # Try to find JSON object in the text
    if not content.startswith('{'):
        match = re.search(r'\{[\s\S]*\}', content)
        if match:
            content = match.group(0)

    data = json.loads(content)
    if not isinstance(data, dict):
        raise ValueError('Model returned non-object JSON')

    # Cap unbounded lists to avoid memory/CPU exhaustion (DoS)
    if 'slides' in data:
        data['slides'] = data['slides'][:MAX_SLIDES]
    if 'sections' in data:
        data['sections'] = data['sections'][:MAX_SECTIONS]
    return data


# ── Markdown escaping (prevent stored XSS in MDX output) ─────────

def _escape_mdx_text(text: str) -> str:
    """Escape text destined for a Markdown slide deck."""
    if not isinstance(text, str):
        text = str(text)
    # Escape HTML tags and common Markdown link injection
    return (text.replace('&', '&amp;').replace('<', '&lt;')
                .replace('>', '&gt;'))


# ── PPTX generation ──────────────────────────────────────────────

def _create_pptx(data: dict, output_path: Path) -> None:
    """Create a PowerPoint file from structured data."""
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.dml.color import RGBColor

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Color scheme
    PRIMARY = RGBColor(0x25, 0x63, 0xEB)    # Blue
    DARK = RGBColor(0x1E, 0x29, 0x3B)       # Dark navy
    LIGHT = RGBColor(0xF8, 0xFA, 0xFC)      # Light gray
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)
    ACCENT = RGBColor(0x10, 0xB9, 0x81)     # Green

    def add_bg(slide, color):
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = color

    # Title slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    add_bg(slide, DARK)

    # Title
    left, top, width, height = Inches(1), Inches(2), Inches(11), Inches(2)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = data.get('title', 'Presentation')
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    if data.get('subtitle'):
        left, top, width, height = Inches(2), Inches(4.2), Inches(9), Inches(1)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = data['subtitle']
        p.font.size = Pt(20)
        p.font.color.rgb = RGBColor(0x94, 0xA3, 0xB8)
        p.alignment = PP_ALIGN.CENTER

    # Content slides
    for i, slide_data in enumerate(data.get('slides', [])):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        add_bg(slide, WHITE)

        # Slide number
        txBox = slide.shapes.add_textbox(Inches(12), Inches(0.3), Inches(1), Inches(0.5))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = str(i + 2)
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(0x94, 0xA3, 0xB8)
        p.alignment = PP_ALIGN.RIGHT

        # Accent bar
        from pptx.util import Emu
        shape = slide.shapes.add_shape(
            1, Inches(0.5), Inches(0.8), Inches(0.08), Inches(0.6)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = PRIMARY
        shape.line.fill.background()

        # Title
        txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.7), Inches(11), Inches(0.8))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = slide_data.get('title', f'Slide {i+2}')
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = DARK

        # Content bullets
        txBox = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11), Inches(5))
        tf = txBox.text_frame
        tf.word_wrap = True
        content_items = slide_data.get('content', [])
        for j, item in enumerate(content_items):
            if j == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = f"  •  {item}"
            p.font.size = Pt(18)
            p.font.color.rgb = RGBColor(0x33, 0x41, 0x55)
            p.space_after = Pt(12)

        # Speaker notes
        if slide_data.get('notes'):
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = slide_data['notes']

    # Thank you slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, PRIMARY)
    txBox = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11), Inches(2))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = 'Thank You! 🙏'
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    prs.save(str(output_path))


# ── DOCX generation ──────────────────────────────────────────────

def _create_docx(data: dict, output_path: Path) -> None:
    """Create a Word document from structured data."""
    from docx import Document
    from docx.shared import Pt, Inches, RGBColor
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.enum.style import WD_STYLE_TYPE

    doc = Document()

    # Style setup
    style = doc.styles['Normal']
    font = style.font
    font.size = Pt(11)
    font.name = 'Arial'

    # Title
    title = doc.add_heading(data.get('title', 'Document'), level=0)
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER

    if data.get('subtitle'):
        subtitle = doc.add_paragraph(data['subtitle'])
        subtitle.alignment = WD_ALIGN_PARAGRAPH.CENTER
        subtitle.runs[0].font.size = Pt(14)
        subtitle.runs[0].font.color.rgb = RGBColor(0x64, 0x74, 0x8B)

    doc.add_paragraph('')  # Spacer

    # Sections
    for section_data in data.get('sections', []):
        doc.add_heading(section_data.get('heading', 'Section'), level=1)

        if section_data.get('content'):
            doc.add_paragraph(section_data['content'])

        for subsection in section_data.get('subsections', []):
            doc.add_heading(subsection.get('heading', ''), level=2)
            if subsection.get('content'):
                doc.add_paragraph(subsection['content'])

    # Footer
    doc.add_paragraph('')
    footer = doc.add_paragraph('Generated by Sanjhubai — sanjhubai.ir')
    footer.alignment = WD_ALIGN_PARAGRAPH.CENTER
    footer.runs[0].font.size = Pt(9)
    footer.runs[0].font.color.rgb = RGBColor(0x94, 0xA3, 0xB8)

    doc.save(str(output_path))


# ── Markdown Slide Deck generation ───────────────────────────────

def _create_mdx(data: dict, output_path: Path) -> None:
    """Create a Markdown slide deck (compatible with Marp/reveal.js)."""
    lines = []

    # Title slide
    lines.append('---')
    lines.append(f'# {data.get("title", "Presentation")}')
    if data.get('subtitle'):
        lines.append(f'\n### {data["subtitle"]}')
    lines.append('\n---\n')

    # Content slides
    for slide_data in data.get('slides', []):
        lines.append('---')
        lines.append(f'\n## {_escape_mdx_text(slide_data.get("title", "Slide"))}\n')
        for item in slide_data.get('content', []):
            lines.append(f'- {_escape_mdx_text(item)}')
        if slide_data.get('notes'):
            lines.append(f'\n<!-- Notes: {_escape_mdx_text(slide_data["notes"])} -->')
        lines.append('')

    # End
    lines.append('---')
    lines.append('\n# Thank You! 🙏')
    lines.append('\n---')

    output_path.write_text('\n'.join(lines), encoding='utf-8')


# ── Endpoints ────────────────────────────────────────────────────

@router.post('/v1/documents/generate')
async def generate_document(request: Request) -> JSONResponse:
    """Generate a document from a prompt."""
    uid = await _get_user_id(request)
    if not uid:
        return JSONResponse({'error': {'message': 'لطفاً وارد حساب خود شوید'}}, status_code=401)

    try:
        body = await request.json()
    except Exception:
        return JSONResponse({'error': {'message': 'درخواست نامعتبر'}}, status_code=400)

    prompt = body.get('prompt', '').strip()
    doc_type = body.get('type', 'pptx').lower()  # pptx, docx, mdx
    model = body.get('model', 'mimo-v2.5-pro')

    if not prompt:
        return JSONResponse({'error': {'message': 'متن درخواست الزامی است'}}, status_code=400)

    if len(prompt) > MAX_PROMPT_LENGTH:
        return JSONResponse(
            {'error': {'message': f'متن درخواست نباید بیشتر از {MAX_PROMPT_LENGTH} کاراکتر باشد'}},
            status_code=413,
        )

    if doc_type not in ('pptx', 'docx', 'mdx'):
        return JSONResponse({'error': {'message': 'نوع فایل نامعتبر. pptx, docx, یا mdx'}}, status_code=400)

    if model not in ALLOWED_MODELS:
        return JSONResponse({'error': {'message': 'مدل انتخابی پشتیبانی نمیشود'}}, status_code=400)

    doc_id = uuid.uuid4().hex[:12]
    ext_map = {'pptx': '.pptx', 'docx': '.docx', 'mdx': '.md'}
    output_path = DOC_STORAGE / f'{doc_id}{ext_map[doc_type]}'

    try:
        # Generate content with AI
        start = time.time()
        data = await _generate_content(prompt, doc_type if doc_type != 'mdx' else 'pptx', model)
        gen_time = time.time() - start

        # Create document
        if doc_type == 'pptx':
            _create_pptx(data, output_path)
        elif doc_type == 'docx':
            _create_docx(data, output_path)
        else:
            _create_mdx(data, output_path)

        file_size = output_path.stat().st_size

        # Register document
        _doc_registry[doc_id] = {
            'id': doc_id,
            'user_id': uid,
            'prompt': prompt[:MAX_PROMPT_LENGTH],
            'type': doc_type,
            'model': model,
            'title': data.get('title', 'Untitled'),
            'filename': f'{doc_id}{ext_map[doc_type]}',
            'file_size': file_size,
            'generation_time': round(gen_time, 1),
            'slides_count': len(data.get('slides', [])),
            'sections_count': len(data.get('sections', [])),
            'created_at': datetime.now(timezone.utc).isoformat(),
        }

        return JSONResponse({
            'id': doc_id,
            'type': doc_type,
            'title': data.get('title', 'Untitled'),
            'filename': f'{doc_id}{ext_map[doc_type]}',
            'file_size': file_size,
            'generation_time': round(gen_time, 1),
            'slides_count': len(data.get('slides', [])),
            'sections_count': len(data.get('sections', [])),
            'download_url': f'/v1/documents/{doc_id}/download',
        })

    except json.JSONDecodeError as e:
        logger.error(f'JSON parse error: {e}')
        return JSONResponse({'error': {'message': 'خطا در تولید محتوا. دوباره تلاش کنید.'}}, status_code=500)
    except Exception as e:
        logger.error(f'Document generation failed: {e}', exc_info=True)
        return JSONResponse({'error': {'message': 'خطا در تولید سند. لطفاً دوباره تلاش کنید.'}}, status_code=500)


@router.get('/v1/documents')
async def list_documents(request: Request) -> JSONResponse:
    """List user's generated documents."""
    uid = await _get_user_id(request)
    if not uid:
        return JSONResponse({'error': {'message': 'لطفاً وارد حساب خود شوید'}}, status_code=401)

    docs = [d for d in _doc_registry.values() if d['user_id'] == uid]
    docs.sort(key=lambda d: d['created_at'], reverse=True)
    return JSONResponse({'documents': docs})


@router.get('/v1/documents/{doc_id}/download')
async def download_document(doc_id: str, request: Request) -> Response:
    """Download a generated document."""
    uid = await _get_user_id(request)
    if not uid:
        return JSONResponse({'error': {'message': 'لطفاً وارد حساب خود شوید'}}, status_code=401)

    # Strict doc_id validation (defense-in-depth vs path traversal)
    if not DOC_ID_RE.match(doc_id or ''):
        return JSONResponse({'error': {'message': 'شناسه سند نامعتبر است'}}, status_code=400)

    doc = _doc_registry.get(doc_id)
    if not doc or doc['user_id'] != uid:
        return JSONResponse({'error': {'message': 'سند یافت نشد'}}, status_code=404)

    ext_map = {'pptx': '.pptx', 'docx': '.docx', 'mdx': '.md'}
    path = DOC_STORAGE / f'{doc_id}{ext_map[doc["type"]]}'

    # Ensure resolved path stays inside DOC_STORAGE (defense-in-depth)
    if not str(path.resolve()).startswith(str(DOC_STORAGE.resolve()) + os.sep):
        return JSONResponse({'error': {'message': 'سند یافت نشد'}}, status_code=404)

    if not path.exists():
        return JSONResponse({'error': {'message': 'فایل یافت نشد'}}, status_code=404)

    mime_map = {
        'pptx': 'application/vnd.openxmlformats-officedocument.presentationml.presentation',
        'docx': 'application/vnd.openxmlformats-officedocument.wordprocessingml.document',
        'mdx': 'text/markdown',
    }

    # Sanitize download filename
    safe_title = re.sub(r'[^A-Za-z0-9_\u0600-\u06FF\- ]', '', doc['title'])[:80] or 'document'

    return FileResponse(
        path=str(path),
        media_type=mime_map.get(doc['type'], 'application/octet-stream'),
        filename=f'{safe_title}{ext_map[doc["type"]]}',
    )


@router.delete('/v1/documents/{doc_id}')
async def delete_document(doc_id: str, request: Request) -> JSONResponse:
    """Delete a generated document."""
    uid = await _get_user_id(request)
    if not uid:
        return JSONResponse({'error': {'message': 'لطفاً وارد حساب خود شوید'}}, status_code=401)

    if not DOC_ID_RE.match(doc_id or ''):
        return JSONResponse({'error': {'message': 'شناسه سند نامعتبر است'}}, status_code=400)

    doc = _doc_registry.get(doc_id)
    if not doc or doc['user_id'] != uid:
        return JSONResponse({'error': {'message': 'سند یافت نشد'}}, status_code=404)

    ext_map = {'pptx': '.pptx', 'docx': '.docx', 'mdx': '.md'}
    path = DOC_STORAGE / f'{doc_id}{ext_map[doc["type"]]}'
    path.unlink(missing_ok=True)
    del _doc_registry[doc_id]

    return JSONResponse({'status': 'deleted'})
